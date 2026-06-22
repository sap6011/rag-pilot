from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
import nbformat


def load_pdf(path: Path):
    """
    One entry per page. For image-heavy pages (sparse text layer), we render
    the page and ask a VLM to caption it, appending the caption to the text.
    """
    from ingest.vlm_captioner import render_pdf_pages, needs_captioning, caption_image

    reader = PdfReader(path)

    # Lazy render: only render pages if any page needs captioning
    rendered_pages = None

    results = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        meta = {"source": path.name, "page": i + 1, "type": "pdf"}

        if needs_captioning(text):
            # First time we need a caption — render all pages now
            if rendered_pages is None:
                print(f"  Rendering {path.name} for VLM captioning...")
                rendered_pages = render_pdf_pages(path)

            print(f"  Captioning page {i + 1}...")
            caption = caption_image(rendered_pages[i])
            if caption:
                text = f"{text}\n\n[Visual content]: {caption}"

        results.append((text, meta))

    return results


def load_pptx(path: Path):
    """One entry per slide."""
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        out.append((text, {"source": path.name, "slide": i + 1, "type": "pptx"}))
    return out


def load_docx(path: Path):
    """One entry for the whole document."""
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs)
    return [(text, {"source": path.name, "type": "docx"})]


def load_ipynb(path: Path):
    """One entry per cell, with cell type in metadata."""
    nb = nbformat.read(path, as_version=4)
    out = []
    for i, cell in enumerate(nb.cells):
        if cell.cell_type in ("markdown", "code"):
            out.append((
                cell.source,
                {"source": path.name, "cell": i, "cell_type": cell.cell_type, "type": "ipynb"},
            ))
    return out


def load_rmd(path: Path):
    """R Markdown is just text."""
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [(text, {"source": path.name, "type": "rmd"})]


LOADERS = {
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    ".docx": load_docx,
    ".ipynb": load_ipynb,
    ".rmd": load_rmd,
}


def load_any(path: Path):
    """Dispatch to the right loader based on file extension."""
    fn = LOADERS.get(path.suffix.lower())
    if not fn:
        return []
    try:
        return fn(path)
    except Exception as e:
        print(f"Failed to load {path.name}: {e}")
        return []